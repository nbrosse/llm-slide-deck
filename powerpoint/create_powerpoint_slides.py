from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

powerpoint_path = Path(__file__).parent.resolve()
repo_path = powerpoint_path.parent.resolve()

# Assets paths
edf_logo_path = str(repo_path / 'images/edf-logo.png')
slide1_bg_path = str(repo_path / 'images/slide-1-background.jpg')
slide3_acteurs_path = str(repo_path / 'images/slide-3-acteurs.jpg')

# --- Helper function to add a footer ---
def add_footer(slide):
    slide.shapes.add_picture(edf_logo_path, Inches(8.5), Inches(6.8), width=Inches(1))
    # Note: Adding page numbers and text is easier with a master slide template.
    # We are adding them manually here for demonstration.
    
# --- Create Presentation ---
prs = Presentation()
# Use a widescreen format (16:9)
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ==============================================================================
# SLIDE 1: Title Slide
# ==============================================================================
slide1_layout = prs.slide_layouts[6] # Blank layout
slide1 = prs.slides.add_slide(slide1_layout)

# Add background image
slide1.shapes.add_picture(slide1_bg_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

# Add semi-transparent overlay
left, top, width, height = Inches(1.5), Inches(1), Inches(7), Inches(5.5)
shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
shape.fill.transparency = 0.25
shape.line.fill.background() # No outline

# Add content
slide1.shapes.add_picture(edf_logo_path, Inches(2), Inches(1.5), width=Inches(1.2))
slide1.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1.5)).text = "Appel d'offres\néolien terrestre"
slide1.shapes.add_textbox(Inches(2), Inches(4.2), Inches(6), Inches(0.5)).text = "(publié par la Commission de Régulation de l'Energie le 28 Avril 2017)"
txBox = slide1.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(1))
p = txBox.text_frame.paragraphs[0]
p.text = "LIVRET D'ACCEUIL\nPRODUCTEUR"
p.font.color.rgb = RGBColor(237, 125, 49) # Orange color
p.font.bold = True
p.font.size = Pt(28)

# ==============================================================================
# SLIDE 2: Table of Contents (SOMMAIRE)
# ==============================================================================
slide2_layout = prs.slide_layouts[6]
slide2 = prs.slides.add_slide(slide2_layout)
slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(1)).text = "SOMMAIRE"

toc_items = [
    ("Préambule", RGBColor(0, 112, 192)),
    ("Présentation des acteurs", RGBColor(0, 32, 96)),
    ("Parcours de contractualisation", RGBColor(84, 130, 53)),
    ("Check-list des démarches", RGBColor(154, 196, 14)),
    ("Questions - Réponses", RGBColor(237, 125, 49)),
    ("Adresses utiles", RGBColor(255, 69, 0))
]

top_pos = Inches(1.5)
for text, color in toc_items:
    shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), top_pos, Inches(4), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.text = text
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape.text_frame.paragraphs[0].font.bold = True
    top_pos += Inches(0.8)

add_footer(slide2)


# ==============================================================================
# SLIDE 3: Preamble
# ==============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

# Orange title bar
shape = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(237, 125, 49)
shape.line.fill.background()
title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(4), Inches(1))
title.text_frame.paragraphs[0].text = "Préambule"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
title.text_frame.paragraphs[0].font.size = Pt(44)

# Main content (from PDF)
content = (
    "Ce document s’adresse uniquement aux lauréats de l’appel d’offres « Installations de production d’électricité à partir de l’énergie mécanique du vent, implantées à terre » (FET17).",
    "Ce document résume, sous une forme simplifiée, les étapes nécessaires à l’élaboration du contrat de complément de rémunération pour une installation lauréate de l’appel d’offres éolien terrestre, lancé par la Commission de Régulation de l’Energie (CRE) le 28 avril 2017.",
    "Dans le cadre des missions de service public prévues par l’article L311-12 du code de l’énergie, EDF est tenue de conclure un contrat de complément de rémunération avec les lauréats retenus à l’issue de l’appel d’offres."
)
txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(4))
tf = txBox.text_frame
tf.clear()
p0 = tf.paragraphs[0]
p0.text = content[0]
p0.level = 0
p0.font.size = Pt(16)
for line in content[1:]:
    p = tf.add_paragraph()
    p.text = line
    p.level = 0
    p.font.size = Pt(16)

# Info boxes
box1 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.5), Inches(3), Inches(2.5))
box1.fill.solid()
box1.fill.fore_color.rgb = RGBColor(0, 32, 96)
box1.text = (
    "Ce livret ne saurait engager la responsabilité d’EDF quant aux obligations du producteur "
    "de s’assurer qu’il respecte le cadre législatif et règlementaire applicable à son installation."
)
box1.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

box2 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(4.2), Inches(3), Inches(3))
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(0, 32, 96)
box2.text = (
    "Le lauréat s’engage à mettre en service et à exploiter une installation en tous points conforme "
    "aux stipulations du cahier des charges de l’appel d’offres et aux caractéristiques décrites dans son offre "
    "(seuls les écarts mentionnés dans l’appel d’offres sont tolérés)."
)
box2.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

add_footer(slide3)

# ==============================================================================
# SLIDE 4: Actors Diagram (Simplified)
# ==============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
# Add orange title bar (as in slide 3)
shape = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(237, 125, 49); shape.line.fill.background()
title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(1))
title.text_frame.paragraphs[0].text = "Présentation des acteurs"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
title.text_frame.paragraphs[0].font.size = Pt(44)

# Central element with image fill
center_x, center_y, radius = Inches(5), Inches(4), Inches(1.5)
producer_shape = slide4.shapes.add_shape(MSO_SHAPE.OVAL, center_x - radius, center_y - radius, radius*2, radius*2)
producer_shape.text = "Producteur"
producer_shape.text_frame.paragraphs[0].font.bold = True
producer_shape.text_frame.paragraphs[0].font.size = Pt(24)
producer_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
producer_shape.fill.background() # Important: remove default fill
# Add picture fill
producer_shape.line.fill.background()
producer_shape.shadow.inherit = False
# This part is tricky; image fill is not directly supported in the same way as the UI.
# A common workaround is to place an image and crop it to a circle shape, which is complex.
# For simplicity, we will place a picture behind a transparent circle.
slide4.shapes.add_picture(slide3_acteurs_path, center_x - radius, center_y - radius, width=radius*2, height=radius*2)
producer_shape.fill.solid()
producer_shape.fill.fore_color.rgb = RGBColor(0,0,0)
producer_shape.fill.transparency = 1.0 # Make it see-through

# Add surrounding text bubbles (simplified positions)
actors = {
    "Commission de Régulation de l’Energie\nPour répondre à l’Appel d’Offres": (Inches(1), Inches(2)),
    "Marché de l’électricité\nPour vendre mon énergie produite.\nSeul un contrat de complément de rémunération est signé avec EDF.": (Inches(7), Inches(2)),
    "Préfet de Région / DGEC\nPour toute modification (d’exploitant, de puissance, …)": (Inches(1), Inches(5)),
    "Gestionnaire du Réseau de Distribution ou de Transport : Enedis, ELD ou RTE\nPour obtenir un contrat d’accès au réseau et mettre en service l’installation.": (Inches(3), Inches(6)),
    "EDF OA (Obligations d’Achat)\nPour obtenir le contrat de complément de rémunération correspondant à l’appel d’offres.": (Inches(6.5), Inches(5))
}
for text, (left, top) in actors.items():
    shape = slide4.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(2.5), Inches(1.5))
    shape.text = text
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_footer(slide4)

# ==============================================================================
# SLIDE 5: Parcours de contractualisation
# ==============================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])

# Orange title bar
shape = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(237, 125, 49)
shape.line.fill.background()
title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
title.text_frame.paragraphs[0].text = "Parcours de contractualisation"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
title.text_frame.paragraphs[0].font.size = Pt(40)

# Steps content (from PDF)
steps = [
    ("1 Demande de raccordement", "J’effectue ma demande de raccordement auprès du gestionnaire de réseau (maximum 2 mois après la désignation)."),
    ("2 Demande de contrat", "Au plus près de l’achèvement de mon installation, j’envoie ma demande de contrat à EDF OA accompagnée des pièces listées page 7."),
    ("3 Notification de la date projetée de prise d’effet", "Je notifie à EDF OA la date projetée de prise d’effet de mon contrat. La notification s’effectue par voie postale ou par voie dématérialisée."),
    ("4 Mise en service du raccordement", "Je prends rendez-vous avec mon gestionnaire de réseau pour mettre en service le raccordement de mon installation au réseau."),
    ("5 Achèvement de l’installation et attestation de conformité", "J’achève mon installation dans un délai de 36 mois à compter de la date de désignation. Je fais établir, par un organisme agréé, une attestation de conformité qui confirmera le respect du cahier des charges de l’appel d’offres éolien terrestre et la conformité de l’installation aux éléments mentionnés dans mon offre de candidature."),
    ("6 Signature du contrat de complément de rémunération", "Dans le cadre du processus de signature, EDF OA m’adresse mon contrat de complément de rémunération."),
    ("7 Facture et règlement", "J’émets mes factures mensuellement, sur la base des données de facturation transmises par le gestionnaire de réseau selon les modalités définies aux conditions générales de mon contrat de complément de rémunération et les transmets à EDF OA. De plus, en début d’année civile, j’adresse à EDF OA la facture ou l’avoir de régularisation annuelle conformément aux dispositions des conditions générales de mon contrat.")
]

content_box = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
tf5 = content_box.text_frame
tf5.clear()

for idx, (step_title, step_text) in enumerate(steps):
    if idx == 0:
        p = tf5.paragraphs[0]
    else:
        p = tf5.add_paragraph()
    p.text = f"{step_title} – {step_text}"
    p.level = 0
    p.font.size = Pt(14)
    p.space_after = Pt(6)

add_footer(slide5)

# --- Save the presentation ---
prs.save(str(powerpoint_path / "EDF_Presentation_Python.pptx"))
print("Presentation 'EDF_Presentation_Python.pptx' created successfully.")