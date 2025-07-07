import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from PIL import Image

# --- Configuration & Constants ---

# File paths
ASSETS_DIR = 'assets'
EDF_LOGO_PATH = os.path.join(ASSETS_DIR, 'edf_logo.png')
SLIDE1_BG_PATH = os.path.join(ASSETS_DIR, 'slide1_bg.jpeg')
SLIDE4_DIAGRAM_PATH = os.path.join(ASSETS_DIR, 'slide4_diagram.png')
OVERLAY_PATH = os.path.join(ASSETS_DIR, 'overlay.png')
OUTPUT_FILENAME = 'EDF_Appel_d_offres_Eolien.pptx'

# Color Palette (RGB)
EDF_ORANGE = RGBColor(255, 102, 0)
EDF_BLUE_DARK = RGBColor(0, 51, 102)
HEADER_ORANGE = RGBColor(245, 124, 0)
TEXT_GRAY = RGBColor(89, 89, 89)
INFO_BOX_BLUE = RGBColor(44, 62, 80)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY_BG = RGBColor(245, 245, 245)

def create_transparent_overlay(path, size, color=(255, 255, 255), opacity=0.85):
    """Creates a semi-transparent image file using Pillow."""
    if not os.path.exists(ASSETS_DIR):
        os.makedirs(ASSETS_DIR)
    img = Image.new('RGBA', size, color + (int(255 * opacity),))
    img.save(path, 'PNG')

def add_slide_header(slide, title_text):
    """Adds the standard orange header bar to a slide."""
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = HEADER_ORANGE
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(8), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = title_text
    p.font.name = 'Open Sans'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

def create_slide_1(prs):
    """Creates the title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # Background image
    slide.shapes.add_picture(SLIDE1_BG_PATH, 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    # Semi-transparent overlay (created on the fly)
    overlay_size_pixels = (900, 600)
    create_transparent_overlay(OVERLAY_PATH, overlay_size_pixels, opacity=0.15) # Invert opacity for PIL
    slide.shapes.add_picture(OVERLAY_PATH, Inches(1.5), Inches(0.8), width=Inches(10.33))
    
    # EDF Logo
    slide.shapes.add_picture(EDF_LOGO_PATH, Inches(2.0), Inches(1.2), height=Inches(0.5))

    # Title
    title_box = slide.shapes.add_textbox(Inches(2.0), Inches(2.0), Inches(9), Inches(1.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Appel d'offres\néolien terrestre"
    p.font.name = 'Open Sans'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = EDF_BLUE_DARK
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(2.0), Inches(3.6), Inches(9), Inches(0.5))
    p = subtitle_box.text_frame.paragraphs[0]
    p.text = "(publié par la Commission de Régulation de l'Energie le 28 Avril 2017)"
    p.font.name = 'Open Sans'
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.LINE, Inches(2.5), Inches(4.5), Inches(8), Inches(0))
    line.line.color.rgb = RGBColor(221, 221, 221)
    line.line.width = Pt(1)

    # Livret
    livret_box = slide.shapes.add_textbox(Inches(2.0), Inches(4.8), Inches(9), Inches(1))
    p = livret_box.text_frame.paragraphs[0]
    p.text = "LIVRET D'ACCUEIL PRODUCTEUR"
    p.font.name = 'Open Sans'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = EDF_ORANGE
    p.alignment = PP_ALIGN.CENTER

def create_slide_2(prs):
    """Creates the 'Sommaire' (Table of Contents) slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    p = title_box.text_frame.paragraphs[0]
    p.text = "SOMMAIRE"
    p.font.name = 'Open Sans'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT_GRAY

    # Sommaire items
    items = [
        ("Préambule", RGBColor(41, 121, 255)),
        ("Présentation des acteurs", RGBColor(0, 45, 98)),
        ("Parcours de contractualisation", RGBColor(85, 139, 47)),
        ("Check-list des démarches", RGBColor(175, 180, 43)),
        ("Questions - Réponses", RGBColor(245, 124, 0)),
        ("Adresses utiles", RGBColor(230, 74, 25)),
    ]
    
    top = Inches(1.8)
    for text, color in items:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(7), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        
        p = shape.text_frame.paragraphs[0]
        p.text = text
        p.font.name = 'Open Sans'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT
        shape.text_frame.margin_left = Inches(0.2)
        top += Inches(0.8)

def create_slide_3(prs):
    """Creates the 'Préambule' slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    add_slide_header(slide, "Préambule")

    # Left column text
    left_col = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7), Inches(5))
    tf = left_col.text_frame
    tf.word_wrap = True
    
    points = [
        "Ce document s'adresse uniquement aux lauréats de l'appel d'offres « Installations de production d'électricité à partir de l'énergie mécanique du vent, implantées à terre » (FET17).",
        "Ce document résume, sous une forme simplifiée, les étapes nécessaires à l'élaboration du contrat de complément de rémunération...",
        "Dans le cadre des missions de service public prévues par l'article L311-12 du code de l'énergie, EDF est tenue de conclure un contrat..."
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.font.name = 'Open Sans'
        p.font.size = Pt(16)
        p.level = 0
        p.space_after = Pt(12)
        # In python-pptx, bullet color follows font color
        p.font.color.rgb = EDF_ORANGE
        # Manually set the text color back
        run = p.add_run()
        run.text = '' # This seems odd but it's a way to reset following text
        p.font.color.rgb = TEXT_GRAY

    # Right column info boxes
    # Box 1
    info1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(1.8), Inches(4.8), Inches(2.0))
    info1.fill.solid()
    info1.fill.fore_color.rgb = INFO_BOX_BLUE
    info1.line.fill.background()
    p = info1.text_frame.paragraphs[0]
    p.text = "   Ce livret ne saurait engager la responsabilité d'EDF quant aux obligations du producteur de s'assurer qu'il respecte le cadre législatif et règlementaire applicable à son installation."
    p.font.name = 'Open Sans'; p.font.size = Pt(12); p.font.color.rgb = WHITE
    info1.text_frame.margin_left = info1.text_frame.margin_right = Inches(0.2)
    
    # "i" icon for box 1
    i_box1 = slide.shapes.add_textbox(Inches(8.1), Inches(1.8), Inches(0.5), Inches(0.5))
    p = i_box1.text_frame.paragraphs[0]
    p.text = 'i'; p.font.name = 'Times New Roman'; p.font.size = Pt(40); p.font.italic = True; p.font.color.rgb = WHITE
    
    # Box 2
    info2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(4.2), Inches(4.8), Inches(2.5))
    info2.fill.solid()
    info2.fill.fore_color.rgb = INFO_BOX_BLUE
    info2.line.fill.background()
    p = info2.text_frame.paragraphs[0]
    p.text = "   Le lauréat s'engage à mettre en service et à exploiter une installation en tous points conforme aux stipulations du cahier des charges..."
    p.font.name = 'Open Sans'; p.font.size = Pt(12); p.font.color.rgb = WHITE
    info2.text_frame.margin_left = info2.text_frame.margin_right = Inches(0.2)
    
    # "i" icon for box 2
    i_box2 = slide.shapes.add_textbox(Inches(8.1), Inches(4.2), Inches(0.5), Inches(0.5))
    p = i_box2.text_frame.paragraphs[0]
    p.text = 'i'; p.font.name = 'Times New Roman'; p.font.size = Pt(40); p.font.italic = True; p.font.color.rgb = WHITE

def create_slide_4(prs):
    """Creates the 'Acteurs' slide with the diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Présentation des acteurs")
    
    # Insert the pre-rendered diagram image
    slide.shapes.add_picture(SLIDE4_DIAGRAM_PATH, Inches(0.6), Inches(1.2), width=Inches(12))

def create_slide_5(prs):
    """Creates the 'Parcours' slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Parcours de contractualisation")
    
    steps = [
        ("Demande de raccordement", "J'effectue ma demande de raccordement auprès du gestionnaire de réseau (maximum 2 mois après la désignation).", RGBColor(25, 118, 210)),
        ("Demande de contrat", "Au plus près de l'achèvement de mon installation, j'envoie ma demande de contrat à EDF OA...", RGBColor(48, 63, 159)),
        ("Notification de la date...", "Je notifie à EDF OA la date projetée de prise d'effet de mon contrat...", RGBColor(251, 192, 45)),
        ("Mise en service...", "Je prends rendez vous avec mon gestionnaire de réseau pour mettre en service...", RGBColor(245, 124, 0)),
        ("Achèvement de l'installation...", "J'achève mon installation dans un délai de 36 mois à compter de la date de désignation...", RGBColor(211, 47, 47)),
        ("Signature du contrat...", "Dans le cadre du processus de signature, EDF OA m'adresse mon contrat de complément de rémunération.", RGBColor(56, 142, 60)),
        ("Facture et règlement", "Je mets mes factures mensuellement, sur la base des données de facturation transmises par le gestionnaire de réseau...", RGBColor(123, 31, 162)),
    ]
    
    top = Inches(1.4)
    row_height = Inches(0.8)
    
    for i, (step_text, detail_text, color) in enumerate(steps):
        # Left colored box
        step_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(4.5), row_height)
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = color
        step_box.line.fill.background()
        tf = step_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        tf.margin_left = Inches(0.1)

        run_num = p.add_run()
        run_num.text = f"{i+1} "
        run_num.font.name = 'Open Sans'
        run_num.font.bold = True
        run_num.font.size = Pt(24)
        run_num.font.color.rgb = WHITE if i != 2 else TEXT_GRAY # Special case for yellow bg

        run_text = p.add_run()
        run_text.text = step_text
        run_text.font.name = 'Open Sans'
        run_text.font.bold = True
        run_text.font.size = Pt(14)
        run_text.font.color.rgb = WHITE if i != 2 else TEXT_GRAY

        # Right detail box
        detail_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.05), top, Inches(7.8), row_height)
        detail_box.fill.solid()
        detail_box.fill.fore_color.rgb = LIGHT_GRAY_BG
        detail_box.line.fill.background()
        
        p = detail_box.text_frame.paragraphs[0]
        p.text = detail_text
        p.font.name = 'Open Sans'
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_GRAY
        p.alignment = PP_ALIGN.LEFT
        detail_box.text_frame.margin_left = Inches(0.15)
        detail_box.text_frame.vertical_anchor = PP_ALIGN.CENTER

        top += row_height + Inches(0.08)

def main():
    """Main function to generate the presentation."""
    prs = Presentation()
    # Set slide size to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    print("Generating Slide 1: Title...")
    create_slide_1(prs)
    
    print("Generating Slide 2: Sommaire...")
    create_slide_2(prs)
    
    print("Generating Slide 3: Préambule...")
    create_slide_3(prs)
    
    print("Generating Slide 4: Acteurs...")
    create_slide_4(prs)
    
    print("Generating Slide 5: Parcours...")
    create_slide_5(prs)
    
    prs.save(OUTPUT_FILENAME)
    print(f"\nPresentation saved as '{OUTPUT_FILENAME}'")

if __name__ == '__main__':
    main()